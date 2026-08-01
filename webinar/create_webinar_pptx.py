"""
Virgin CPR Webinar — PowerPoint Generator
Run: pip install python-pptx  (if not installed)
Then: python3 create_webinar_pptx.py
Output: virgin_cpr_webinar.pptx (in the same folder)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ──────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
GOLD         = RGBColor(0xFF, 0xC3, 0x00)   # gold
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xCC, 0xCC)
GREEN        = RGBColor(0x00, 0xD4, 0x7E)
RED          = RGBColor(0xFF, 0x4D, 0x4D)
ACCENT_BLUE  = RGBColor(0x4A, 0xA8, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # truly blank layout


# ── Helper functions ─────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_divider(slide, top, color=GOLD):
    add_rect(slide, 0.5, top, 12.33, 0.04, color)


def title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, BG_DARK)
    # gold accent bar on left
    add_rect(slide, 0, 0, 0.18, 7.5, GOLD)
    # title
    add_textbox(slide, title,
                0.4, 2.2, 12.5, 1.5,
                font_size=40, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    # divider
    add_divider(slide, 3.85)
    # subtitle
    add_textbox(slide, subtitle,
                0.4, 4.0, 12.5, 1.0,
                font_size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    return slide


def section_slide(prs, section_num, section_title, subtitle=""):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, BG_DARK)
    add_rect(slide, 0, 3.4, 13.33, 0.7, GOLD)
    add_textbox(slide, f"HOUR {section_num}",
                0.5, 0.6, 12, 1.0,
                font_size=16, bold=True, color=GOLD)
    add_textbox(slide, section_title,
                0.5, 1.2, 12, 1.8,
                font_size=36, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, subtitle,
                    0.5, 3.5, 12, 0.6,
                    font_size=18, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    return slide


def content_slide(prs, title, bullets, note=""):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, BG_DARK)
    # header bar
    add_rect(slide, 0, 0, 13.33, 1.1, RGBColor(0x12, 0x27, 0x3A))
    add_textbox(slide, title,
                0.4, 0.15, 12.5, 0.8,
                font_size=24, bold=True, color=GOLD)
    add_divider(slide, 1.15)

    # bullets
    y = 1.35
    for bullet in bullets:
        if bullet.startswith("##"):          # sub-heading
            add_textbox(slide, bullet[2:].strip(),
                        0.5, y, 12, 0.45,
                        font_size=16, bold=True, color=ACCENT_BLUE)
            y += 0.5
        elif bullet.startswith("✅") or bullet.startswith("⭐"):
            add_textbox(slide, bullet,
                        0.5, y, 12, 0.4,
                        font_size=15, color=GREEN)
            y += 0.45
        elif bullet.startswith("❌") or bullet.startswith("⚠️"):
            add_textbox(slide, bullet,
                        0.5, y, 12, 0.4,
                        font_size=15, color=RED)
            y += 0.45
        else:
            add_textbox(slide, f"  {bullet}",
                        0.5, y, 12, 0.4,
                        font_size=15, color=WHITE)
            y += 0.43

    if note:
        add_rect(slide, 0.4, 6.7, 12.53, 0.55, RGBColor(0x1A, 0x35, 0x50))
        add_textbox(slide, f"💡  {note}",
                    0.6, 6.72, 12.1, 0.45,
                    font_size=13, italic=True, color=GOLD)
    return slide


def formula_slide(prs):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, BG_DARK)
    add_rect(slide, 0, 0, 13.33, 1.1, RGBColor(0x12, 0x27, 0x3A))
    add_textbox(slide, "The CPR Formula — Never Forget This",
                0.4, 0.15, 12.5, 0.8,
                font_size=24, bold=True, color=GOLD)
    add_divider(slide, 1.15)

    formulas = [
        ("Pivot  =", "(PDH + PDL + PDC)  ÷  3"),
        ("BC     =", "(PDH + PDL)  ÷  2"),
        ("TC     =", "2 × Pivot  −  BC"),
        ("Band   =", "[ min(TC, BC)  →  max(TC, BC) ]"),
        ("Width  =", "TCPR  −  BCPR"),
    ]
    y = 1.5
    for label, formula in formulas:
        add_rect(slide, 0.5, y, 3.8, 0.55, RGBColor(0x1A, 0x35, 0x50))
        add_textbox(slide, label,
                    0.6, y + 0.05, 3.6, 0.45,
                    font_size=16, bold=True, color=ACCENT_BLUE)
        add_rect(slide, 4.5, y, 8.3, 0.55, RGBColor(0x0A, 0x14, 0x1F))
        add_textbox(slide, formula,
                    4.6, y + 0.05, 8.1, 0.45,
                    font_size=16, color=WHITE, bold=True)
        y += 0.72

    add_rect(slide, 0.4, 6.3, 12.53, 0.75, RGBColor(0xFF, 0xC3, 0x00, ))
    add_textbox(slide,
                "VIRGIN  =  The band has NEVER been touched since it formed.",
                0.6, 6.35, 12.1, 0.6,
                font_size=17, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
    return slide


def pattern_slide(prs, pattern_id, name, pair, key_candle, entry_rule, result):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, BG_DARK)
    add_rect(slide, 0, 0, 13.33, 1.1, RGBColor(0x12, 0x27, 0x3A))
    add_textbox(slide, f"Pattern {pattern_id} — {name}",
                0.4, 0.12, 12.5, 0.85,
                font_size=26, bold=True, color=GOLD)
    add_divider(slide, 1.15)

    # pair badge
    add_rect(slide, 0.5, 1.3, 2.2, 0.5, ACCENT_BLUE)
    add_textbox(slide, pair,
                0.55, 1.32, 2.1, 0.44,
                font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # key candle box
    add_rect(slide, 0.5, 2.0, 12.33, 0.55, RGBColor(0x1A, 0x35, 0x50))
    add_textbox(slide, f"Key Candle:  {key_candle}",
                0.65, 2.03, 12.0, 0.48,
                font_size=15, color=ACCENT_BLUE, bold=True)

    # entry rule
    add_rect(slide, 0.5, 2.75, 12.33, 0.55, RGBColor(0x00, 0x33, 0x1A))
    add_textbox(slide, f"Entry Rule:  {entry_rule}",
                0.65, 2.78, 12.0, 0.48,
                font_size=15, color=GREEN, bold=True)

    # result
    add_rect(slide, 0.5, 3.5, 12.33, 0.55, RGBColor(0x33, 0x1A, 0x00))
    add_textbox(slide, f"Result:  {result}",
                0.65, 3.53, 12.0, 0.48,
                font_size=15, color=RGBColor(0xFF, 0xA5, 0x00), bold=True)

    # chart placeholder
    add_rect(slide, 0.5, 4.25, 12.33, 2.75, RGBColor(0x0A, 0x14, 0x1F),
             line_color=RGBColor(0x33, 0x55, 0x77))
    add_textbox(slide, "[ Paste your TradingView chart screenshot here ]",
                0.5, 5.3, 12.33, 0.6,
                font_size=14, italic=True, color=RGBColor(0x55, 0x77, 0x99),
                align=PP_ALIGN.CENTER)
    return slide


def rules_table_slide(prs):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, BG_DARK)
    add_rect(slide, 0, 0, 13.33, 1.1, RGBColor(0x12, 0x27, 0x3A))
    add_textbox(slide, "Trade Entry Checklist — Every Single Trade",
                0.4, 0.12, 12.5, 0.85,
                font_size=24, bold=True, color=GOLD)
    add_divider(slide, 1.15)

    rows = [
        ("Session",      "London open (12:30 IST) or NY open (18:30 IST)",    "❌ Never Asian"),
        ("Band Status",  "Must still be VIRGIN — zero touches",                "❌ Touched = skip"),
        ("Band Age",     "Older band = stronger reaction (3-day > 1-day)",     "✅ Note the date"),
        ("Entry candle", "Rejection pin bar OR engulfing close outside band",  "❌ No inside-bar entry"),
        ("Stop Loss",    "5 pips above TCPR (sell) or below BCPR (buy)",       "✅ Fixed rule"),
        ("Position Size","Risk max 1% of account per trade",                   "✅ Calculate before entry"),
        ("News",         "No trades 15 min before/after high-impact news",     "❌ Skip news candles"),
    ]

    y = 1.3
    headers = ["Checklist Item", "Rule", "Note"]
    col_w = [2.5, 6.8, 3.5]
    col_x = [0.4, 2.95, 9.8]
    h_colors = [GOLD, GOLD, GOLD]

    # header row
    for i, h in enumerate(headers):
        add_rect(slide, col_x[i], y, col_w[i], 0.45, RGBColor(0x1A, 0x35, 0x50))
        add_textbox(slide, h, col_x[i]+0.05, y+0.04, col_w[i]-0.1, 0.38,
                    font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    y += 0.48

    for item, rule, note in rows:
        row_bg = RGBColor(0x0D, 0x1B, 0x2A) if (rows.index((item,rule,note)) % 2 == 0) else RGBColor(0x12, 0x27, 0x3A)
        for i, txt in enumerate([item, rule, note]):
            add_rect(slide, col_x[i], y, col_w[i], 0.44, row_bg,
                     line_color=RGBColor(0x1A, 0x35, 0x50))
            c = RED if txt.startswith("❌") else (GREEN if txt.startswith("✅") else WHITE)
            add_textbox(slide, txt, col_x[i]+0.05, y+0.04,
                        col_w[i]-0.1, 0.38, font_size=12, color=c)
        y += 0.47
    return slide


def closing_slide(prs):
    slide = prs.slides.add_slide(BLANK)
    set_bg(slide, BG_DARK)
    add_rect(slide, 0, 0, 0.18, 7.5, GOLD)
    add_rect(slide, 0, 5.8, 13.33, 1.7, RGBColor(0x0A, 0x14, 0x1F))

    add_textbox(slide, "Your Next 7 Days",
                0.4, 0.6, 12.5, 1.0,
                font_size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_divider(slide, 1.75)

    steps = [
        "Day 1–2  →  Calculate CPR manually for XAUUSD, EURUSD, GBPUSD every morning",
        "Day 3–4  →  Mark virgin bands on your chart — paper trade, no real money",
        "Day 5–6  →  Wait for Pattern A or B only — take ONE trade setup per day",
        "Day 7    →  Review your journal — did you follow the checklist every time?",
    ]
    y = 2.0
    for step in steps:
        add_rect(slide, 0.5, y, 12.33, 0.52, RGBColor(0x12, 0x27, 0x3A))
        add_textbox(slide, step, 0.65, y+0.06, 12.0, 0.42,
                    font_size=15, color=WHITE)
        y += 0.62

    add_textbox(slide,
                "\"A Virgin CPR is a promise the market has not kept yet — and it always keeps its promises.\"",
                0.4, 5.85, 12.5, 0.65,
                font_size=13, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Thank You — Questions?",
                0.4, 6.6, 12.5, 0.55,
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    return slide


# ── Build deck ───────────────────────────────────────────────────────────────

# Slide 1 — Title
title_slide(prs,
    "Virgin CPR — Forex Mastery Workshop",
    "5-Hour Live Webinar  |  From Zero to Execution  |  ₹4999")

# Slide 2 — Agenda
content_slide(prs, "What You Will Learn Today (Agenda)", [
    "Hour 1  →  Why CPR works — the math and the magnet concept",
    "Hour 2  →  The 5 proven patterns — real charts, candle by candle",
    "Hour 3  →  Live chart demo — mark up today's CPR together",
    "Hour 4  →  Trade rules, SL placement, position sizing",
    "Hour 5  →  Trade simulation quiz + backtest proof + open Q&A",
    "──────────────────────────────────────────",
    "✅ You will leave with a complete entry checklist",
    "✅ You will leave with a 30-day practice journal template",
    "✅ Recording shared within 24 hours",
], note="Keep your TradingView chart open throughout the session")

# Slide 3 — Why This Webinar
content_slide(prs, "Why This Webinar? — Not Like the Others", [
    "## What most forex webinars give you:",
    "❌ 10+ strategies → you are confused, you don't master any one",
    "❌ ₹10,000–₹50,000 fees → no bot, no tools, no follow-up",
    "❌ Theory-heavy → no real chart walkthroughs, no entry rules",
    "──────────────────────────────────────────",
    "## What you get here:",
    "✅ ONE strategy — Virgin CPR — mastered completely in 5 hours",
    "✅ Covers 12 currency pairs every single day — systematic, not guesswork",
    "✅ A working Telegram Bot that scans all 12 pairs and sends you alerts — FREE",
    "✅ Real charts, real numbers, real entry/SL/TP — nothing hidden",
    "──────────────────────────────────────────",
    "→  You don't need 10 strategies. You need ONE strategy done right.",
], note="\"One method, executed perfectly on 12 pairs, beats 10 half-understood methods.\"")

# Slide 4 — 12 pairs + Free Bot
content_slide(prs, "12 Pairs. 1 Strategy. Every Day. + Free Bot", [
    "## The 12 pairs we scan every morning:",
    "XAUUSD  |  EURUSD  |  GBPUSD  |  USDJPY  |  AUDUSD  |  USDCAD",
    "USDCHF  |  NZDUSD  |  EURJPY  |  GBPJPY  |  EURGBP  |  GBPAUD",
    "──────────────────────────────────────────",
    "## Manual scan takes only 15 minutes every morning.",
    "→  For each pair: calculate CPR, check virgin status, note band width",
    "→  Shortlist the 2–3 best setups to watch during London open",
    "──────────────────────────────────────────",
    "## FREE Telegram Bot (included with this webinar):",
    "✅ Bot scans all 12 pairs automatically every morning",
    "✅ Sends you a Telegram message: which pairs have active Virgin CPR today",
    "✅ Saves you 15 minutes every single day",
    "✅ Bot code is open — you can see exactly how it works",
], note="The bot does not tell you to trade — it tells you WHERE to look. You make the decision.")

# Slide 5 — Indian Trader Trap
content_slide(prs, "The Indian Trader's Trap — Nifty & BankNifty Options", [
    "## Why most Indian traders lose money on options:",
    "❌ Time Decay (Theta) — your option LOSES value every single day, even if price doesn't move",
    "❌ Expiry pressure — every Thursday, options go to zero if out of money",
    "❌ Premium crush after news — IV drops, your option loses value even if direction was right",
    "❌ Market makers know your expiry date — they pin price to hurt max retail holders",
    "──────────────────────────────────────────",
    "## Real example:",
    "→  You buy BankNifty 51000 CE at ₹200. Nifty goes sideways for 2 days.",
    "→  Your option is now worth ₹80 — you lost 60% even though you were not wrong.",
    "→  This is Theta eating your premium. You cannot stop it.",
    "──────────────────────────────────────────",
    "⚠️  In options, TIME is your enemy. The clock is always running against buyers.",
], note="This is not a criticism of options. It is a fact of how options are priced.")

# Slide 6 — Why Forex
content_slide(prs, "Why Forex? — No Time Decay. No Expiry. Full Flexibility.", [
    "## What Forex gives you that options don't:",
    "✅ NO time decay — a trade does not lose value just by sitting there",
    "✅ NO expiry — you hold a trade 1 hour or 1 week, your choice",
    "✅ NO theta, NO gamma, NO vega — pure price action, no Greek complexity",
    "✅ Market open 24 hours, 5 days a week — trade around your job",
    "──────────────────────────────────────────",
    "## Comparison:",
    "→  Options: Buy at ₹200, sit 2 days → ₹80 (lost ₹120 without moving)",
    "→  Forex:   Buy XAUUSD, sit 2 days → same price = same P&L (no decay)",
    "──────────────────────────────────────────",
    "## For salaried Indians especially:",
    "✅ London open is 12:30 PM IST — lunch break trade",
    "✅ NY open is 6:30 PM IST — after office trade",
    "✅ No need to sit in front of screen all day",
], note="You don't need to quit your job to trade forex. Two sessions. One hour each. That is enough.")

# Slide 7 — Funded Account Challenge
content_slide(prs, "Start Smart — Funded Account Challenge First", [
    "## What is a Funded Account (Prop Firm)?",
    "→  You pay a one-time fee (₹5,000–₹15,000) to attempt a trading challenge",
    "→  You pass the challenge using the rules you learn today",
    "→  The firm gives you ₹5 Lakh – ₹40 Lakh of THEIR money to trade",
    "→  You keep 80–90% of all profits. You risk ZERO of your savings.",
    "──────────────────────────────────────────",
    "## Examples of prop firms:",
    "✅ FTMO  |  FundedNext  |  The5ers  |  MyFundedFX",
    "──────────────────────────────────────────",
    "## The smart path for a new trader:",
    "Step 1  →  Learn Virgin CPR strategy (today)",
    "Step 2  →  Paper trade for 30 days — prove it to yourself",
    "Step 3  →  Attempt a funded account challenge with small fee",
    "Step 4  →  Pass → trade firm's money → keep profits",
    "Step 5  →  Only THEN consider putting your own money at risk",
    "──────────────────────────────────────────",
    "✅ This is the safest way to go from learner to paid trader.",
], note="Many of my students passed their FTMO challenge within 60 days of learning this strategy.")

# Slide 8 — Prerequisites
content_slide(prs, "Before We Begin — Are You Ready?", [
    "## You need to know these basics:",
    "→  Read a candlestick chart (bullish, bearish, pin bar, engulfing)",
    "→  What is Support and Resistance (swing high / swing low)",
    "→  What is Stop Loss (SL) and Take Profit (TP)",
    "→  You have TradingView or MT5 open RIGHT NOW",
    "→  You know Asian / London / NY session times in IST",
    "──────────────────────────────────────────",
    "❌ If you don't know these — bookmark the recording and study the Glossary PDF first",
], note="Glossary PDF has been shared in the WhatsApp group")

# Slide 4 — HOUR 1 section
section_slide(prs, "1", "Why CPR Works", "The Math + The Magnet Concept")

# Slide 5 — CPR Formula
formula_slide(prs)

# Slide 6 — What is Virgin
content_slide(prs, "What Makes a CPR 'Virgin'?", [
    "Virgin = The band has NEVER been touched since it formed",
    "Touch condition:   candle_high ≥ BCPR   AND   candle_low ≤ TCPR",
    "──────────────────────────────────────────",
    "## Why it matters:",
    "→  Institutions place orders at CPR zones at the START of the day",
    "→  If price never visits that zone → those orders are STILL sitting there",
    "→  When price finally arrives → the reaction is sharp and fast",
    "──────────────────────────────────────────",
    "## Band age rule:",
    "✅ 3-day-old virgin band  →  VERY strong reaction expected",
    "✅ 1-day-old virgin band  →  moderate reaction",
], note="\"A Virgin CPR is a promise the market has not kept yet — it always keeps its promises.\"")

# Slide 7 — Band Width
content_slide(prs, "Band Width — Tells You What Kind of Day to Expect", [
    "## Narrow CPR (width < ATR)",
    "✅ Market is likely to TREND strongly today",
    "✅ Price approaches band → sharp rejection or clean breakout",
    "✅ Best day for Pattern A (Approach & Snap) and Pattern B (False Break)",
    "──────────────────────────────────────────",
    "## Wide CPR (width > ATR)",
    "→  Market is likely to RANGE or consolidate",
    "→  Price may chop inside the band — harder to trade",
    "→  Best for Pattern C (Band as Launchpad) with order flow confirmation",
    "──────────────────────────────────────────",
    "⚠️  Always check band width BEFORE the session — not during",
], note="Narrow CPR = trending day expected. Wide CPR = be careful, use order flow filter.")

# Slide 8 — HOUR 2 section
section_slide(prs, "2", "The 5 Patterns", "Real Charts · Candle by Candle · Entry & Exit")

# Slide 9–13 — Patterns
pattern_slide(prs, "A", "Approach and Snap",
    "XAUUSD · M15 · 24-Feb-2026",
    "13:15 IST — Bearish pin bar. Wick enters band, close BELOW BCPR.",
    "Sell at close of rejection candle. SL = 5 pts above TCPR. Target = PDL area.",
    "Price fell from 2919.20 → 2905.00 by 16:00 IST.  +14.2 pts = 2R ✅")

pattern_slide(prs, "B", "False Break + Trap",
    "EURUSD · M15 · 25-Feb-2026",
    "12:45 — Wick ABOVE TCPR, but candle CLOSED back inside band (retail trapped).",
    "Wait for next candle to close BELOW BCPR. Sell at that close.",
    "Price fell from 1.0470 → 1.0442.  +28 pips = 1.4R ✅")

pattern_slide(prs, "C", "Band as Launchpad",
    "GBPUSD · M15 · 26-Feb-2026",
    "13:00 — Strong bullish candle. Close ABOVE TCPR after order flow showed ACCUMULATION.",
    "Buy at close above TCPR. SL = below mid-band − 5 pips.",
    "Price rallied from 1.2601 → 1.2648.  +47 pips = 2.9R ✅")

pattern_slide(prs, "D", "Stacked Bands Compression",
    "USDJPY · M15 · 27-Feb-2026",
    "Two virgin bands only 6.5 pips apart. Price chops between them 2.5 hours.",
    "NO TRADE inside compression zone. Enter only when price exits BELOW lower band.",
    "Sell at 152.190. Price fell to 151.820.  +37 pips = 1.9R ✅")

pattern_slide(prs, "E", "Morning Gap Into Band",
    "XAUUSD · M15 · 28-Feb-2026",
    "Overnight gap. Price opened BELOW the virgin band. Band is 7.4 pts above open.",
    "DO NOT CHASE the gap. Wait for price to approach band and show rejection/breakout.",
    "Patience play — wait for London session to bring price to the band.")

# Slide 14 — HOUR 3 section
section_slide(prs, "3", "Live Chart Demo", "Mark Today's CPR Together")

content_slide(prs, "Live Chart — Step by Step (Do This With Me)", [
    "Step 1  →  Open TradingView. Go to your pair (XAUUSD or EURUSD).",
    "Step 2  →  Switch to Daily chart. Note: PDH, PDL, PDC from yesterday.",
    "Step 3  →  Calculate Pivot, BC, TC using the formula (or use the indicator).",
    "Step 4  →  Draw a rectangle on the band — mark BCPR and TCPR.",
    "Step 5  →  Check: has ANY candle since market open touched this band?",
    "──────────────────────────────────────────",
    "✅ If NO touch found → you have a LIVE Virgin CPR right now",
    "✅ Note the band width — narrow or wide?",
    "✅ Note the band age — how many days old?",
    "──────────────────────────────────────────",
    "⚠️  Now watch it. Don't touch it until London open.",
], note="Every student should have their own chart marked by end of this hour")

# Slide 15 — HOUR 4 section
section_slide(prs, "4", "Trade Rules & Management", "Entry · SL · Position Size · Avoid List")

# Slide 16 — Rules table
rules_table_slide(prs)

# Slide 17 — Position sizing
content_slide(prs, "Position Sizing — Protect Your Account First", [
    "## The 1% Rule:",
    "Risk = Account Balance × 1%",
    "Example: ₹50,000 account → Risk per trade = ₹500",
    "──────────────────────────────────────────",
    "## Lot Size Formula:",
    "Lot Size  =  Risk Amount  ÷  (SL in pips × Pip Value)",
    "Example (EURUSD, SL = 20 pips, pip value = ₹85):",
    "Lot Size  =  500  ÷  (20 × 85)  =  0.29 lots",
    "──────────────────────────────────────────",
    "❌ Never increase lot size to recover a loss",
    "❌ Never move SL further away after entry",
    "✅ Calculate BEFORE you click Buy/Sell — not after",
], note="Use the position size calculator in the shared Excel sheet")

# Slide 18 — When NOT to trade
content_slide(prs, "When NOT to Trade — Save Yourself First", [
    "❌ During Asian session (05:00–12:30 IST) — low volume, fake moves",
    "❌ Inside a stacked bands compression zone — price will chop",
    "❌ 15 minutes before / after high-impact news (NFP, CPI, Fed rate)",
    "❌ When the band has already been touched today",
    "❌ When band width is extremely wide (>3× ATR) — no clean reaction",
    "❌ On Fridays after 18:30 IST — weekend gap risk",
    "──────────────────────────────────────────",
    "## The hardest skill in trading:",
    "✅ Doing nothing is a valid trade decision",
    "✅ Preserve capital → live to trade another day",
], note="Most losses come from trading when there is NO setup. Patience is the strategy.")

# Slide 19 — HOUR 5 section
section_slide(prs, "5", "Simulation + Proof + Q&A", "Test Your Knowledge · Backtest Results")

# Slide 20 — Quiz format
content_slide(prs, "Trade Simulation Quiz — You Decide First", [
    "I will show you 3 historical charts.",
    "For each chart, answer BEFORE I reveal the outcome:",
    "──────────────────────────────────────────",
    "## Question 1:",
    "→  Is the CPR virgin?",
    "→  What pattern do you see?",
    "→  Do you enter? If yes — Buy or Sell?",
    "→  Where is your SL?",
    "──────────────────────────────────────────",
    "✅ Compare your answer with the actual trade outcome",
    "✅ Scoring: 3/3 = you are ready. 2/3 = one more week of paper trading. 1/3 = replay Hour 2.",
], note="This is how you know if you have truly understood the strategy — not just watched it")

# Slide 21 — Backtest summary
content_slide(prs, "Backtest Results — The Proof", [
    "## Feb 2026 — 5-Day Study (5 pairs)::",
    "Total setups identified  :  18",
    "Valid virgin setups      :  11",
    "Trades taken (per rules) :  8",
    "Winners                  :  6   (75% win rate)",
    "Average R on winners     :  2.1R",
    "Average R on losers      :  −1.0R",
    "──────────────────────────────────────────",
    "## Net result over 5 days:  +8.6R",
    "✅ On ₹50,000 account at 1% risk  →  +₹4,300 in 5 trading days",
    "──────────────────────────────────────────",
    "⚠️  Past results do not guarantee future performance.",
    "⚠️  These numbers are from rule-based backtesting, not live trading.",
], note="The edge comes from FOLLOWING THE RULES — not from every single trade winning")

# Slide 22 — Closing / 7 days plan
closing_slide(prs)

# ── Save ─────────────────────────────────────────────────────────────────────
output_file = "virgin_cpr_webinar.pptx"
prs.save(output_file)
print(f"\n✅  Saved → {output_file}")
print("   Open it in PowerPoint or Google Slides.")
print("   Replace chart placeholder boxes with your actual TradingView screenshots.")
